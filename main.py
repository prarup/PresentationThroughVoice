from pptx import Presentation
from pptx.util import Inches
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE
import speech_recognition as sr
import pyttsx3
import os
import time
import webbrowser
from pynput.keyboard import Key, Controller
kb = Controller()

print("Speech Recognition System for Teaching Assistance")
engine = pyttsx3.init()
pr1 = Presentation()

def respond(audio):
    engine.setProperty("rate", 200)
    engine.say(audio)
    engine.runAndWait()

def createaprep():
    respond('presentation is being created')
    respond('adding a demo slide')
    slide1_register = pr1.slide_layouts[0]
    slide1 = pr1.slides.add_slide(slide1_register)
    title1 = slide1.shapes.title
    subtitle1 = slide1.placeholders[1]
    title1.text = 'demo1'
    subtitle1.text = 'subs like comment ring the bell icon'
    pr1.save('demo.pptx')
    path = "C:\\Users\\Prarup\\PycharmProjects\\pptThroughVoice\\demo.pptx"
    os.startfile(f"{path}")
    i =1
    while i<100:
            addanother = SpeechCommand()
            #addaslide = presentationcommand()
            if 'add a slide' in addanother:
                os.system('TASKKILL /F /IM powerpnt.exe')
                time.sleep(1)
                slide2_register = pr1.slide_layouts[1]
                slide2 = pr1.slides.add_slide(slide2_register)
                title2 = slide2.shapes.title
                title2.text = 'Demo2'
                bullet_point_box = slide2.shapes
                bullet_points_lvl1 = bullet_point_box.placeholders[1]
                bullet_points_lvl1.text = 'subs'

                pr1.save('demo.pptx')
                path = "C:\\Users\\Prarup\\PycharmProjects\\pptThroughVoice\\demo.pptx"
                os.startfile(f"{path}")
                i += 1
                time.sleep(2)

            if 'add image' in addanother:
                os.system('TASKKILL /F /IM powerpnt.exe')
                time.sleep(1)
                slide3_register = pr1.slide_layouts[5]
                slide3 = pr1.slides.add_slide(slide3_register)
                title3 = slide3.shapes.title
                title3.text = 'Pics'
                img1 = 'cars.jpg'
                from_left = Inches(3)
                from_top = Inches(2)
                add_picture = slide3.shapes.add_picture(img1, from_left, from_top)
                pr1.save('demo.pptx')
                path = "C:\\Users\\Prarup\\PycharmProjects\\pptThroughVoice\\demo.pptx"
                os.startfile(f"{path}")
                i += 1
                time.sleep(2)

            if 'add chart' in addanother:
                os.system('TASKKILL /F /IM powerpnt.exe')
                time.sleep(1)
                slide4_register = pr1.slide_layouts[5]
                slide4 = pr1.slides.add_slide(slide4_register)
                title4 = slide4.shapes.title
                title4.text = 'graphs'
                graphinfo = CategoryChartData()
                graphinfo.categories = ['a', 'b', 'c']
                graphinfo.add_series('series1', (15, 11, 18))
                left_graph = Inches(3)
                top_graph = Inches(2)
                width_graph = Inches(6)
                height_graph = Inches(5)
                graph1_frame = slide4.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, left_graph, top_graph, width_graph,
                                                       height_graph, graphinfo)
                graph1 = graph1_frame.chart
                pr1.save('demo.pptx')
                path = "C:\\Users\\Prarup\\PycharmProjects\\pptThroughVoice\\demo.pptx"
                os.startfile(f"{path}")
                i += 1
                time.sleep(2)

            if 'add table' in addanother:
                os.system('TASKKILL /F /IM powerpnt.exe')
                time.sleep(1)
                slide5_register = pr1.slide_layouts[5]
                slide5 = pr1.slides.add_slide(slide5_register)
                title5 = slide5.shapes.title
                title5.text = 'table'
                left_table = Inches(3)
                top_table = Inches(2)
                width_table = Inches(6)
                height_table = Inches(5)
                table1_frame = slide5.shapes.add_table(3, 4, left_table, top_table, width_table,
                                                       height_table)
                table1 = table1_frame.table
                pr1.save('demo.pptx')
                path = "C:\\Users\\Prarup\\PycharmProjects\\pptThroughVoice\\demo.pptx"
                os.startfile(f"{path}")
                i += 1
                time.sleep(2)

            if 'cancel' in addanother:
                os.system('TASKKILL /F /IM powerpnt.exe')

def SpeechCommand(ask=False):
    r = sr.Recognizer()
    with sr.Microphone() as source:
        if ask:
            print(ask)
        respond("How can I help You?")
        print('Listening...')
        audio = r.listen(source)
        try:
            speechData = r.recognize_google(audio, language='en-in')
            print(speechData)
        except Exception as e:
            print(e)
            respond("Say again Please")
            SpeechCommand()
        if 'End' in speechData:
            exit(0)
        return speechData
def search():
    respond("what do you want me to search")
    find = SpeechCommand()
    url = 'http://google.com/search?q=' + str(find)
    webbrowser.get().open(url)
    respond("here is what i found for" + str(find))

def location():
    respond("what location do you want me to find")
    location = SpeechCommand()
    url = 'http://google.nl/maps/place/' + str(location) + '/&amp;'
    webbrowser.get().open(url)
    respond("here is what i found for" + str(location))

def weather():
    respond("the weather for today is as follows")
    url = 'http://google.com/search?q=weather'
    webbrowser.get().open(url)

def sharemarket():
    respond("the share market for today is as follows")
    url = 'www.nepalstock.com'
    webbrowser.get().open(url)

def playpre():
    respond('What subject do you want to present')
    chpre = SpeechCommand()
    if 'embedded system' in chpre:
        respond('please select chapter')
        chap = SpeechCommand()
        if 'first chapter' in chap:
            path = "C:\\Users\\Prarup\\PycharmProjects\\pptThroughVoice\\chap1.pptx"
            os.startfile(f"{path}")
            time.sleep(2)

        if 'second chapter' in chap:
            path = "C:\\Users\\Prarup\\PycharmProjects\\pptThroughVoice\\chap2.pptx"
            os.startfile(f"{path}")
            time.sleep(2)
    if 'simulation and modelling' in chpre:
        respond("please select chapter")
        chap = SpeechCommand()
        if 'first chapter' in chap:
            path = "C:\\Users\\Prarup\\PycharmProjects\\pptThroughVoice\\chap3.pptx"
            os.startfile(f"{path}")
            time.sleep(2)

        if 'second chapter' in chap:
            path = "C:\\Users\\Prarup\\PycharmProjects\\pptThroughVoice\\chap4.pptx"
            os.startfile(f"{path}")
            time.sleep(2)

    if 'play demo' in chpre:
        path = "C:\\Users\\Prarup\\PycharmProjects\\pptThroughVoice\\demo.pptx"
        os.startfile(f"{path}")
        time.sleep(2)
def keypress():
    time.sleep(1)
    kb.press(Key.f5)
    kb.release(Key.f5)
    j = 1
    while j < 100:
        pprep = SpeechCommand()
        if 'next slide' in pprep:
            kb.press(Key.right)
            kb.release(Key.right)
            j += 1
            time.sleep(2)

        if 'previous slide' in pprep:
            kb.press(Key.left)
            kb.release(Key.left)
            j += 1
            time.sleep(2)

        if 'end presentation' in pprep:
            kb.press(Key.esc)
            kb.release(Key.esc)
            j += 1


        if 'go to start' in pprep:
            kb.press(Key.home)
            kb.release(Key.home)
            j += 1
            time.sleep(2)

        if 'end' in pprep:
            j += 1
            os.system('TASKKILL /F /IM powerpnt.exe')
            break

def response(speechData):
    if 'are you up' in speechData:
        respond("for you, always sir")
    if 'create a presentation' in speechData:
        respond(createaprep())
    if 'search' in speechData:
        respond(search())
    if 'location' in speechData:
        respond(location())
    if 'what is the weather today' in speechData:
        respond(weather())
    if 'share market' in speechData:
        respond(sharemarket())
    if 'play presentation' in speechData:
        respond(playpre())
        respond('Presentation is opening')
        respond(keypress())


speechData = SpeechCommand()
response(speechData)